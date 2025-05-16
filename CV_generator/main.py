import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import os
import math
import re

# === Leitura do Excel ===
df = pd.read_excel("resultado.xlsx")
df.columns = df.columns.str.strip()

# === Funções auxiliares ===
def is_valid(val):
    return val and not (isinstance(val, float) and math.isnan(val))

def extract_skills(row):
    parts = []
    for col in ["Skill", "Specialization Skills", "Specialization  Branch Skills"]:
        if is_valid(row.get(col)):
            parts.append(str(row[col]))
    return "\n".join(parts)

def extract_experience_blocks(description):
    if not is_valid(description):
        return ""
    text = str(description)
    blocks = re.split(r"(?=(?:[A-Z][a-z]{2} \d{4}|\d{4})\s*–\s*(?:[A-Z][a-z]{2} \d{4}|\d{4}))", text)
    results = []
    for i in range(1, len(blocks), 2):
        bloco = (blocks[i] + blocks[i+1]) if i+1 < len(blocks) else blocks[i]
        results.append(bloco.strip())
    def get_year(b):
        m = re.search(r"(\d{4})\s*–", b)
        return int(m.group(1)) if m else 0
    results = sorted(results, key=get_year, reverse=True)
    return "\n\n".join(results[:5])

def add_title(slide, text, top):
    left = Inches(0.5)
    width = Inches(9)
    height = Inches(0.5)
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = Pt(18)
    run.font.bold = True
    run.font.color.rgb = RGBColor(102, 0, 204)  # Roxo
    return top + Inches(0.4)

def add_content(slide, text, top):
    if not is_valid(text):
        return top
    left = Inches(0.5)
    width = Inches(9)
    height = Inches(1.5)
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    for line in str(text).split("\n"):
        p = tf.add_paragraph()
        run = p.add_run()
        run.text = line.strip()
        run.font.size = Pt(14)
    return top + Inches(1.0)

# === Gerar slides por colaborador ===
for _, row in df.iterrows():
    name = row["Worker Name"]
    title = row.get("Job Title", "")
    profile = row.get("Profile", "")
    education = row.get("Education", "")
    experience = extract_experience_blocks(row.get("Description", ""))
    skills = extract_skills(row)
    industries = row.get("Industry Networks", "")
    if is_valid(industries):
        skills += f"\n\nIndustries:\n{industries}"
    photo = row.get("Photo", "").strip()

    # Criar apresentação
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Topo: Nome + Cargo
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    tf = title_box.text_frame
    name_par = tf.paragraphs[0]
    name_par.text = name
    name_par.font.size = Pt(26)
    name_par.font.bold = True

    role_par = tf.add_paragraph()
    role_par.text = title
    role_par.font.size = Pt(18)

    # Inserir foto se existir
    if is_valid(photo) and os.path.exists(photo):
        slide.shapes.add_picture(photo, Inches(8), Inches(0.2), height=Inches(1.5))

    # Corpo: secções com títulos e conteúdo
    top = Inches(1.2)
    top = add_title(slide, "Profile", top)
    top = add_content(slide, profile, top)

    top = add_title(slide, "Education", top)
    top = add_content(slide, education, top)

    top = add_title(slide, "Relevant Skills & Qualifications", top)
    top = add_content(slide, skills, top)

    top = add_title(slide, "Relevant Experience", top)
    top = add_content(slide, experience, top)

    # Criar pasta e guardar
    safe_name = name.replace(" ", "_")
    folder = os.path.join("CVs_Gerados", safe_name)
    os.makedirs(folder, exist_ok=True)
    output_path = os.path.join(folder, f"{safe_name}_CV.pptx")
    prs.save(output_path)
    print(f"✅ Gerado: {output_path}")

print("\n🎯 Todos os CVs foram criados em pastas individuais com layout construído do zero.")
