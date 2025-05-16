# ⬇️ Instala bibliotecas necessárias
!pip install -q pandas python-pptx pillow

# 📁 Upload do Excel
from google.colab import files
uploaded = files.upload()

# 📚 Imports
import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import io
import os
import requests

# Mapeamento meses 3 letras em PT
month_map = {
    "01": "Jan", "02": "Fev", "03": "Mar", "04": "Abr", "05": "Mai", "06": "Jun",
    "07": "Jul", "08": "Ago", "09": "Set", "10": "Out", "11": "Nov", "12": "Dez"
}

def abrevia_mes(data_str):
    try:
        if pd.isna(data_str):
            return ""
        parts = data_str.strip().split('/')
        if len(parts) == 2:
            mm, yyyy = parts
            return f"{month_map.get(mm, mm)} {yyyy}"
        return data_str
    except:
        return data_str

# Ler Excel
filename = list(uploaded.keys())[0]
df = pd.read_excel(io.BytesIO(uploaded[filename]))

def get_start_date(date_range):
    try:
        start = date_range.split("-")[0].strip()
        if '/' in start:
            return pd.to_datetime(start, format="%m/%Y", errors='coerce')
        return pd.NaT
    except:
        return pd.NaT

df['start_date'] = df['Project Date'].apply(get_start_date)

grouped = df.groupby("Worker Name")

os.makedirs("CVs", exist_ok=True)

TITLE_FONT_SIZE = Pt(28)
TEXT_FONT_SIZE = Pt(14)
ACCENT_COLOR = RGBColor(107, 28, 183)

for name, data in grouped:
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    job_title = data["Job Title"].iloc[0]
    profile = data["Profile"].iloc[0] if pd.notna(data["Profile"].iloc[0]) else ""
    
    # Educação (trocar ";" por parágrafos)
    education_list = data["Education"].dropna().unique()
    edu_lines = []
    for edu in education_list:
        parts = str(edu).split(";")
        for part in parts:
            sub_parts = [p.strip() for p in part.split('-')]
            edu_lines.append(" - ".join(sub_parts))
    
    # Skills (máximo 10)
    skills_cols = ["Skill", "Specialization Skills", "Specialization  Branch Skills"]
    skills_set = []
    seen_skills = set()
    for col in skills_cols:
        if col in data.columns:
            for val in data[col].dropna().unique():
                val = str(val).strip()
                if val and val not in seen_skills:
                    skills_set.append(val)
                    seen_skills.add(val)
                    if len(skills_set) >= 10:
                        break
        if len(skills_set) >= 10:
            break
    skills_text = "\n".join(skills_set)
    
    # Industries
    industries = "\n".join(data["Industry Networks"].dropna().unique()) if "Industry Networks" in data.columns else ""
    
    # Languages
    lang_str = data["Language"].iloc[0] if "Language" in data.columns else ""
    languages_list = [x.strip() for x in str(lang_str).split(",")] if lang_str else []
    languages = "\n".join(languages_list)
    
    # Título (nome + job title)
    title_box = slide.shapes.add_textbox(Inches(2), Inches(0.3), Inches(6), Inches(1))
    frame = title_box.text_frame
    p = frame.paragraphs[0]
    p.text = f"{name}\n{job_title}"
    p.runs[0].font.size = TITLE_FONT_SIZE
    p.runs[0].font.color.rgb = ACCENT_COLOR
    
    # Foto
    photo_url = data["Photo"].iloc[0] if "Photo" in data.columns else None
    if isinstance(photo_url, str) and photo_url.startswith("http"):
        try:
            img_data = requests.get(photo_url).content
            with open("temp_photo.png", "wb") as f:
                f.write(img_data)
            slide.shapes.add_picture("temp_photo.png", Inches(0.3), Inches(0.3), Inches(1.5), Inches(1.5))
        except:
            pass
    
    # Profile (caixa larga)
    profile_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.6), Inches(8.5), Inches(1.2))
    profile_frame = profile_box.text_frame
    profile_frame.word_wrap = True
    profile_frame.text = "Profile"
    profile_frame.paragraphs[0].runs[0].font.bold = True
    profile_frame.paragraphs[0].runs[0].font.color.rgb = ACCENT_COLOR
    p = profile_frame.add_paragraph()
    p.text = profile
    p.font.size = TEXT_FONT_SIZE
    
    # Educação
    edu_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.9), Inches(4), Inches(1.5))
    edu_frame = edu_box.text_frame
    edu_frame.word_wrap = True
    edu_frame.text = "Education"
    edu_frame.paragraphs[0].runs[0].font.bold = True
    edu_frame.paragraphs[0].runs[0].font.color.rgb = ACCENT_COLOR
    for line in edu_lines:
        p = edu_frame.add_paragraph()
        p.text = line
        p.font.size = TEXT_FONT_SIZE
    
    # Skills
    skills_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.0), Inches(3), Inches(2))
    skills_frame = skills_box.text_frame
    skills_frame.word_wrap = True
    skills_frame.text = "Skills"
    skills_frame.paragraphs[0].runs[0].font.bold = True
    skills_frame.paragraphs[0].runs[0].font.color.rgb = ACCENT_COLOR
    p = skills_frame.add_paragraph()
    p.text = skills_text
    p.font.size = TEXT_FONT_SIZE
    
    # Industries
    if industries.strip():
        industries_box = slide.shapes.add_textbox(Inches(3.6), Inches(4.0), Inches(3), Inches(1.2))
        industries_frame = industries_box.text_frame
        industries_frame.word_wrap = True
        industries_frame.text = "Industries"
        industries_frame.paragraphs[0].runs[0].font.bold = True
        industries_frame.paragraphs[0].runs[0].font.color.rgb = ACCENT_COLOR
        p = industries_frame.add_paragraph()
        p.text = industries
        p.font.size = TEXT_FONT_SIZE

        # Languages embaixo das industries, se existir
        if languages.strip():
            languages_box = slide.shapes.add_textbox(Inches(3.6), Inches(5.3), Inches(3), Inches(1))
            languages_frame = languages_box.text_frame
            languages_frame.word_wrap = True
            languages_frame.text = "Languages"
            languages_frame.paragraphs[0].runs[0].font.bold = True
            languages_frame.paragraphs[0].runs[0].font.color.rgb = ACCENT_COLOR
            p = languages_frame.add_paragraph()
            p.text = languages
            p.font.size = TEXT_FONT_SIZE
    else:
        # Se não tiver industries, mostrar languages na posição do industries
        if languages.strip():
            languages_box = slide.shapes.add_textbox(Inches(3.6), Inches(4.0), Inches(3), Inches(1))
            languages_frame = languages_box.text_frame
            languages_frame.word_wrap = True
            languages_frame.text = "Languages"
            languages_frame.paragraphs[0].runs[0].font.bold = True
            languages_frame.paragraphs[0].runs[0].font.color.rgb = ACCENT_COLOR
            p = languages_frame.add_paragraph()
            p.text = languages
            p.font.size = TEXT_FONT_SIZE

    # Relevant Experience (mover para a esquerda de 6.8 para 6.0)
    exp_box_height = 6
    exp_box_top = 1.6
    n_proj = min(len(data), 5)
    approx_para_height = 0.3
    total_text_height = (1 + n_proj * 3) * approx_para_height
    top_pos = exp_box_top + max(0, (exp_box_height - total_text_height) / 2)

    exp_box = slide.shapes.add_textbox(Inches(6.0), Inches(top_pos), Inches(3.5), Inches(exp_box_height))
    exp_frame = exp_box.text_frame
    exp_frame.word_wrap = True
    exp_frame.text = "Relevant Experience"
    exp_frame.paragraphs[0].runs[0].font.bold = True
    exp_frame.paragraphs[0].runs[0].font.color.rgb = ACCENT_COLOR
    exp_frame.paragraphs[0].font.size = TEXT_FONT_SIZE

    recent_projects = data.sort_values(by="start_date", ascending=False).head(5)
    for _, row in recent_projects.iterrows():
        proj_date = row["Project Date"]
        proj_name = row["Project Name"]
        desc = row["Description"]

        if isinstance(proj_date, str) and '-' in proj_date:
            try:
                start, end = proj_date.split('-')
                start_fmt = abrevia_mes(start.strip())
                end_fmt = abrevia_mes(end.strip())
                date_str = f"{start_fmt}-----{end_fmt}"
            except:
                date_str = proj_date
        else:
            date_str = proj_date

        for part in [str(date_str), str(proj_name), str(desc)]:
            p = exp_frame.add_paragraph()
            p.text = part.strip()
            p.font.size = TEXT_FONT_SIZE

    # Salvar
    safe_name = name.replace(" ", "_").lower()
    prs.save(f"CVs/{safe_name}.pptx")

# Compactar e download
!zip -r CVs.zip CVs
files.download("CVs.zip")
